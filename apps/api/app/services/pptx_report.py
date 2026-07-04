"""
PPTX Report Service
====================
Generates a PowerPoint presentation from analysis data using python-pptx.
Mirrors the data shape used by pdf_report.generate_pdf / docx_report.generate_docx.
"""

from __future__ import annotations

import io
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

ANAC_BLUE = RGBColor(0x00, 0x3A, 0x70)
DARK      = RGBColor(0x1e, 0x29, 0x3b)
GREY      = RGBColor(0x64, 0x74, 0x8b)
WHITE     = RGBColor(0xff, 0xff, 0xff)

_TIPO_LABELS = {"ciclos": "Ciclos de Fiscalização", "pdf": "Documento PDF", "generic": "Planilha Genérica"}
_ALERT_PREFIX = {"error": "● CRÍTICO", "warning": "▲ ATENÇÃO", "info": "ℹ INFO"}


def _add_title_slide(prs: Presentation, filename: str, created_at: datetime | None, analysis_type: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ANAC Data Insight"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ANAC_BLUE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Briefing Executivo"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = filename
    p3.font.size = Pt(14)
    p3.font.color.rgb = DARK
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    data_str = created_at.strftime("%d/%m/%Y %H:%M") if created_at else "—"
    p4.text = f"{_TIPO_LABELS.get(analysis_type, analysis_type)} · {data_str}"
    p4.font.size = Pt(11)
    p4.font.color.rgb = GREY
    p4.alignment = PP_ALIGN.CENTER


def _add_heading_slide(prs: Presentation, title: str) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ANAC_BLUE
    return slide


def _add_table(slide, rows: list[tuple[str, str]], headers: tuple[str, str] = ("Indicador", "Valor"),
               top: float = 1.2) -> None:
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, 2, Inches(0.5), Inches(top), Inches(9), Inches(0.4 * n_rows))
    table = table_shape.table
    table.columns[0].width = Inches(5.5)
    table.columns[1].width = Inches(3.5)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ANAC_BLUE

    for r, (label, value) in enumerate(rows, start=1):
        table.cell(r, 0).text = label
        table.cell(r, 1).text = str(value)
        for c in range(2):
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(12)


def _add_bullet_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    if not lines:
        return
    slide = _add_heading_slide(prs, title)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK


def generate_pptx(
    *,
    filename: str,
    created_at: datetime | None,
    total_rows: int,
    indicators: dict[str, Any] | None,
    alerts: list[dict] | None,
    ai_summary: dict[str, Any] | None,
    analysis_type: str = "ciclos",
    briefing_extra: dict[str, Any] | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Capa ────────────────────────────────────────────────────────────
    _add_title_slide(prs, filename, created_at, analysis_type)

    ind = indicators or {}

    # ── Sumário executivo / KPIs ─────────────────────────────────────────
    if analysis_type == "pdf":
        slide = _add_heading_slide(prs, "Informações do Documento")
        _add_table(slide, [
            ("Páginas", ind.get("pages", 0)),
            ("Palavras", f"{ind.get('word_count', 0):,}".replace(",", ".")),
            ("Registros", f"{total_rows:,}".replace(",", ".")),
        ])
    else:
        slide = _add_heading_slide(prs, "KPIs Principais")
        _add_table(slide, [
            ("Total de atividades", ind.get("total_atividades", 0)),
            ("Realizadas", ind.get("realizadas", 0)),
            ("Agendadas", ind.get("agendadas", 0)),
            ("Sem agendamento", ind.get("sem_agendamento", 0)),
            ("Taxa de execução (%)", ind.get("taxa_execucao", 0)),
            ("Taxa de agendamento (%)", ind.get("taxa_agendamento", 0)),
            ("Pendências críticas", ind.get("pendencias_criticas", 0)),
        ])

        # ── Ranking por tipo de ciclo ────────────────────────────────────
        by_type = ind.get("by_type")
        if by_type:
            tipo_labels = {"CICLO_BASE": "Ciclo Base", "CICLO_DESEMPENHO": "Desempenho", "NAO_PROGRAMADA": "Não Programadas"}
            rows = [
                (tipo_labels.get(tipo, tipo), f"{vals.get('realizadas', 0)}/{vals.get('total_atividades', 0)} ({vals.get('taxa_execucao', 0):.1f}%)")
                for tipo, vals in by_type.items()
            ]
            slide2 = _add_heading_slide(prs, "Resultado por Tipo de Ciclo")
            _add_table(slide2, rows, headers=("Tipo", "Realizadas / Total"))

    # ── PTA Mensal e comparação histórica (briefing executivo) ───────────
    extra = briefing_extra or {}
    pta_status = extra.get("pta_status")
    if pta_status:
        label = {"realizado": "Realizadas", "agendado": "Agendadas", "sem-agendamento": "Sem agendamento"}
        slide_pta = _add_heading_slide(prs, "PTA Mensal — Situação Atual")
        _add_table(slide_pta, [(label.get(k, k), v) for k, v in pta_status.items()])

    comparison = extra.get("comparison")
    if comparison and comparison.get("average_execution_rate_previous") is not None:
        delta = comparison.get("delta")
        seta = "▲" if (delta or 0) >= 0 else "▼"
        _add_bullet_slide(prs, "Comparação com Ano Anterior", [
            f"Taxa de execução atual: {ind.get('taxa_execucao', 0)}%",
            f"Taxa de execução no mesmo período do ano anterior: {comparison['average_execution_rate_previous']}%",
            f"Variação: {seta} {abs(delta) if delta is not None else 0} p.p.",
        ])

    gerencias_atencao = extra.get("gerencias_atencao")
    if gerencias_atencao:
        _add_bullet_slide(prs, "Gerências em Atenção", [
            f"{g['gerencia']} — {g['criticas']} pendência(s) crítica(s)" for g in gerencias_atencao
        ])

    cidades_atencao = extra.get("cidades_atencao")
    if cidades_atencao:
        _add_bullet_slide(prs, "Cidades em Atenção", [
            f"{c['cidade']} — {c['criticas']} pendência(s) crítica(s)" for c in cidades_atencao
        ])

    monthly_chart = extra.get("monthly_chart")
    if monthly_chart:
        slide_chart = _add_heading_slide(prs, f"Gráfico Mensal — {monthly_chart['mes_label']}")
        chart_data = CategoryChartData()
        chart_data.categories = [monthly_chart["mes_label"]]
        chart_data.add_series("Planejado", (monthly_chart["planejado"],))
        chart_data.add_series("Realizado", (monthly_chart["realizado"],))
        chart_data.add_series("Agendado", (monthly_chart["agendado"],))
        gframe = slide_chart.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.3), Inches(8), Inches(5), chart_data,
        )
        gframe.chart.has_legend = True
        gframe.chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        gframe.chart.legend.include_in_layout = False

    # ── Pendências críticas / Alertas ────────────────────────────────────
    if alerts:
        lines = [
            f"{_ALERT_PREFIX.get(a.get('type', 'info'), '')} {a.get('category', '')} — {a.get('message', '')} ({a.get('count', 0)})"
            for a in alerts
        ]
        _add_bullet_slide(prs, "Pendências Críticas e Alertas", lines)

    # ── Recomendações (IA) ────────────────────────────────────────────────
    if ai_summary:
        resumo = ai_summary.get("resumo_executivo")
        if resumo:
            _add_bullet_slide(prs, "Sumário Executivo", [resumo])
        _add_bullet_slide(prs, "Principais Achados", ai_summary.get("principais_achados") or [])
        _add_bullet_slide(prs, "Riscos Operacionais", ai_summary.get("riscos_operacionais") or [])
        _add_bullet_slide(prs, "Recomendações", ai_summary.get("recomendacoes") or [])

    # ── Metodologia ───────────────────────────────────────────────────────
    _add_bullet_slide(prs, "Metodologia", [
        f"Arquivo: {filename}",
        f"Registros analisados: {total_rows:,}".replace(",", "."),
        "Indicadores calculados automaticamente pela plataforma ANAC Data Insight.",
        "Nenhum dado bruto foi enviado a provedores de IA — apenas indicadores agregados.",
    ])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
