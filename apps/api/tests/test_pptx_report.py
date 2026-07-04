"""Tests for app/services/pptx_report.py."""
import io
from datetime import datetime, timezone

from pptx import Presentation

from app.services.pptx_report import generate_pptx


def test_generate_pptx_ciclos_returns_valid_file():
    data = generate_pptx(
        filename="ciclos.xlsx",
        created_at=datetime.now(timezone.utc),
        total_rows=10,
        indicators={
            "total_atividades": 10, "realizadas": 6, "agendadas": 2, "sem_agendamento": 2,
            "taxa_execucao": 60.0, "taxa_agendamento": 80.0, "pendencias_criticas": 3,
            "by_type": {"CICLO_BASE": {"total_atividades": 5, "realizadas": 3, "taxa_execucao": 60.0}},
        },
        alerts=[{"type": "error", "category": "Sem GIASO", "message": "3 atividades sem GIASO", "count": 3}],
        ai_summary={"resumo_executivo": "Resumo.", "principais_achados": ["Achado 1"],
                    "riscos_operacionais": ["Risco 1"], "recomendacoes": ["Recomendação 1"]},
        analysis_type="ciclos",
    )
    assert isinstance(data, bytes)
    assert len(data) > 0
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 5  # capa, KPIs, by_type, alertas, ai (4 slides) + metodologia


def test_generate_pptx_pdf_type_returns_valid_file():
    data = generate_pptx(
        filename="doc.pdf",
        created_at=datetime.now(timezone.utc),
        total_rows=5,
        indicators={"pages": 5, "word_count": 1200},
        alerts=None,
        ai_summary=None,
        analysis_type="pdf",
    )
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 2  # capa + info do documento + metodologia


def test_generate_pptx_minimal_inputs():
    data = generate_pptx(
        filename="vazio.csv",
        created_at=None,
        total_rows=0,
        indicators=None,
        alerts=None,
        ai_summary=None,
        analysis_type="generic",
    )
    assert isinstance(data, bytes)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 1
